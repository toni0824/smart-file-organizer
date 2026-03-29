import argparse
import json
import os
import re
import shutil
import subprocess
import unicodedata
import urllib.error
import urllib.request
import zipfile
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


OLLAMA_API_URL = "http://localhost:11434/api/chat"
DEFAULT_MODEL = "qwen2.5:1.5b"
DEFAULT_MIN_CONFIDENCE = 85
DEFAULT_MAX_TEXT_CHARS = 6000
CACHE_VERSION = "v6"
CACHE_FILE = Path(__file__).with_name(".path_ai_cache.json")


def find_existing_folder(*names: str) -> Path:
    for name in names:
        candidate = Path.home() / name
        if candidate.exists():
            return candidate
    return Path.home() / names[0]


BASE_FOLDER = find_existing_folder("Downloads", "Transferencias", "Transferências")
DEST_ROOT = Path.home() / "Documents" / "Tecnico"
UNSORTED_FOLDER = DEST_ROOT / "Por_Organizar"

# A ordem importa: regras mais especificas devem aparecer primeiro.
SUBJECT_RULES = {
    "PE": [
        "metodoavaliacaope",
        "formulariope",
        "pe_projeto",
        "resolucao-een",
        "resolucao-eer",
        "een",
        "eer",
    ],
    "RCI": [
        "rci",
        "apresentacaorci",
        "autoavaliacaorci",
        "pautarci",
        "rci_owr",
        "rci manual",
        "rci_manual",
        "rci_s1",
        "rede sobreposta",
        "redesobrepostaencaminhamento",
        "clientserver tcp udp",
        "guiautilizacaoselect",
    ],
    "Telecomunicacoes": [
        "tele_",
        "2o_projeto_de_tele",
        "tele_guia",
        "tele_trabalho",
        "ftele",
        "ist_tel",
        "usrp",
        "adalm-pluto",
        "template-ber",
        "tx_rx_sync",
        "desmodulador",
        "fm_",
        "heterodinagem",
        "detecao coerente",
        "modulacao",
        "apr.04",
        "apr.05",
        "apr.06",
    ],
    "MO": [
        "tclab",
        "heat transfer plant",
        "modelling_of_a_heat_transfer_plant",
        "alpha_tau",
        "openloop_data",
        "t_e_ts",
    ],
    "selec": [
        "iot_module",
        "python-cheat-sheet",
    ],
    "SA": [
        "saut",
        "lidar",
        "odom",
    ],
    "AED": [
        "aed",
    ],
    "MSim": [
        "msim",
        "medicao de grandezas eletricas",
        "condicionamento de sinal",
        "analisador de espetros",
        "sensores, aquisicao e processamento de sinais",
        "waveforms - example",
        "dc measurements",
        "ac measurements",
        "lab1 - feedback",
        "lab2 - proximity detector",
        "lab3",
        "tclab",
        "matlab_daq",
        "aquisicoes",
    ],
    "Sinais e Sistemas": [
        "sinais e sistemas",
        "signals",
        "digitization",
        "oppenheim",
        "ss_0",
        "ss2526",
    ],
    "Fisica": [
        "fisica",
        "fis3",
        "map1",
        "map2",
        "fraccoes simples",
        "formulario",
    ],
    "Programacao Concorrente": [
        "pconc",
    ],
    "Analise de Circuitos": [
        "ac2526",
        "ac2526_map45",
        "ac2526_pauta_lab",
        "ac2526_pratica_p1",
        "ac2526_calendario",
        "kirchhoff",
        "circuitos lineares",
        "regime transitorio",
        "metodos sistematicos",
        "conceitos_basicos",
    ],
    "Arquitetura de Computadores": [
        "arquitetura de computadores",
        "computer architecture",
        "assembly",
        "mips",
    ],
    "Controlo de Voo": [
        "controlo de voo",
        "ils",
        "simulink",
        "root locus",
    ],
    "Ensaios de Voo": [
        "ensaios",
        "ensaios de voo",
        "flight test",
        "propulsivo",
        "estruturais",
    ],
    "PIC": [
        "pic",
        "projeto integrador",
        "propulsao",
        "propulsão",
        "launcher",
        "detritos",
        "aerogel",
    ],
    "Arquitetura": [
        "introducao a arquitetura",
        "introdução à arquitetura",
        "basilica",
        "basílica",
        "estrela",
        "largo do rato",
    ],
    "Propulsao": [
        "propulsao",
        "propulsão",
        "hybrid rocket",
        "nozzle",
        "tanques",
    ],
    "Termica": [
        "tcal",
        "heat transfer",
        "incropera",
        "mass transfer",
        "termica",
    ],
    "Matematica": [
        "calculo",
        "algebra",
        "álgebra",
        "matematica",
        "matemática",
    ],
}

SUBJECT_DESTINATIONS = {
    "PE": DEST_ROOT / "2 ano" / "PE",
    "RCI": DEST_ROOT / "3 ano" / "RCI",
    "Telecomunicacoes": DEST_ROOT / "3 ano" / "tele",
    "MO": DEST_ROOT / "3 ano" / "MO",
    "selec": DEST_ROOT / "3 ano" / "selec",
    "SA": DEST_ROOT / "3 ano" / "SA",
    "AED": DEST_ROOT / "2 ano" / "AED",
    "MSim": DEST_ROOT / "2 ano" / "IM",
    "Sinais e Sistemas": DEST_ROOT / "2 ano" / "SS",
    "Fisica": DEST_ROOT / "2 ano" / "F3",
    "Programacao Concorrente": DEST_ROOT / "2 ano" / "PC",
    "Analise de Circuitos": DEST_ROOT / "2 ano" / "AC",
    "Controlo de Voo": DEST_ROOT / "2 ano" / "Controlo",
}

VALID_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls",
    ".txt", ".zip", ".py", ".m", ".jpg", ".jpeg", ".png",
    ".dmg", ".pkg", ".html", ".csv", ".mov", ".slx", ".grc",
    ".mat", ".dat", ".apk", ".iso",
}

AI_REVIEWABLE_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".txt", ".py", ".m", ".html", ".csv",
}

SKIP_DIRECTORY_SUFFIXES = {
    ".app",
    ".photoslibrary",
}

SKIP_DIRECTORY_NAMES = {
    "__pycache__",
    "node_modules",
    ".git",
    ".venv",
    "venv",
}

ALWAYS_UNSORTED_KEYWORDS = [
    "ticket",
    "tickets",
    "cv",
    "curriculum",
    "offerletter",
    "declaracao",
    "declaracao_de_matricula",
    "declaracao_pt",
    "matricula",
    "screenshot",
    "captura de ecra",
    "captura de tela",
    "whatsapp image",
    "screen recording",
]

ACADEMIC_HINTS = [
    "exame",
    "enunciado",
    "resolucao",
    "solucao",
    "formulario",
    "serie",
    "lab",
    "guia",
    "manual",
    "projeto",
    "trabalho",
    "aula",
    "introducao",
    "notas",
    "map",
    "feedback",
    "template",
    "oscilogramas",
]

def normalize(text: str) -> str:
    text = unicodedata.normalize("NFKD", text)
    text = "".join(char for char in text if not unicodedata.combining(char))
    return text.lower().strip()


def should_force_unsorted(filename: str) -> bool:
    name = normalize(filename)
    for keyword in ALWAYS_UNSORTED_KEYWORDS:
        if keyword in name:
            return True
    return False


def has_academic_hint(text: str) -> bool:
    normalized = normalize(text)
    if not normalized:
        return False

    for keyword in ACADEMIC_HINTS:
        if keyword in normalized:
            return True

    for subject, keywords in SUBJECT_RULES.items():
        if subject == "Matematica":
            continue
        for keyword in keywords:
            if keyword_matches(normalized, keyword):
                return True

    return False


def keyword_matches(name: str, keyword: str) -> bool:
    normalized_keyword = normalize(keyword)

    if " " in normalized_keyword or "_" in normalized_keyword or "-" in normalized_keyword:
        return normalized_keyword in name

    tokens = re.findall(r"[a-z0-9]+", name)
    if len(normalized_keyword) <= 3:
        return normalized_keyword in tokens

    return normalized_keyword in name


def ensure_folder(folder: Path) -> None:
    folder.mkdir(parents=True, exist_ok=True)


def safe_move(file_path: Path, destination_folder: Path, dry_run: bool = False) -> Path:
    if not dry_run:
        ensure_folder(destination_folder)
    destination = destination_folder / file_path.name

    if not destination.exists():
        if not dry_run:
            shutil.move(str(file_path), str(destination))
        return destination

    stem = file_path.stem
    suffix = file_path.suffix
    counter = 1

    while True:
        new_destination = destination_folder / f"{stem}_{counter}{suffix}"
        if not new_destination.exists():
            if not dry_run:
                shutil.move(str(file_path), str(new_destination))
            return new_destination
        counter += 1


def list_directory_hints(dir_path: Path, max_items: int = 15) -> str:
    try:
        children = sorted(dir_path.iterdir(), key=lambda path: path.name.lower())
    except OSError:
        return ""

    names = []
    for child in children[:max_items]:
        names.append(child.name)
    return " | ".join(names)


def collect_directory_preview(dir_path: Path, max_chars: int, max_files: int = 6) -> str:
    parts = [f"[DIR] {dir_path.name}"]
    listing = list_directory_hints(dir_path, max_items=20)
    if listing:
        parts.append(f"[FILES] {listing}")

    added = 0
    try:
        children = sorted(dir_path.rglob("*"), key=lambda path: str(path).lower())
    except OSError:
        return "\n".join(parts)[:max_chars]

    for child in children:
        if added >= max_files:
            break
        if not child.is_file():
            continue
        if child.name.startswith("."):
            continue
        preview = extract_text_preview(child, max(800, max_chars // 3))
        if not preview:
            continue
        parts.append(f"[FILE] {child.name}\n{preview}")
        added += 1
        if len("\n".join(parts)) >= max_chars:
            break

    return "\n".join(parts)[:max_chars]


def should_skip_directory(dir_path: Path) -> bool:
    return dir_path.name in SKIP_DIRECTORY_NAMES or dir_path.suffix.lower() in SKIP_DIRECTORY_SUFFIXES


def destination_for_subject(subject: str | None) -> Path | None:
    if not subject:
        return None
    return SUBJECT_DESTINATIONS.get(subject)


def display_destination(path: Path) -> str:
    try:
        return str(path.relative_to(DEST_ROOT))
    except ValueError:
        return path.name


def classify_with_rules(name_or_hint: str) -> str | None:
    if should_force_unsorted(name_or_hint):
        return None

    name = normalize(name_or_hint)

    for subject, keywords in SUBJECT_RULES.items():
        for keyword in keywords:
            if keyword_matches(name, keyword):
                return subject

    return None


def load_cache() -> dict:
    if not CACHE_FILE.exists():
        return {}

    try:
        return json.loads(CACHE_FILE.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        return {}


def save_cache(cache: dict) -> None:
    CACHE_FILE.write_text(json.dumps(cache, indent=2, ensure_ascii=False), encoding="utf-8")


def build_cache_key(file_path: Path, model: str, aggressive: bool = False) -> str:
    stat = file_path.stat()
    mode = "aggressive" if aggressive else "normal"
    return f"{CACHE_VERSION}:{mode}:{file_path}:{stat.st_size}:{stat.st_mtime_ns}:{model}"


def extract_response_text(response_data: dict) -> str:
    message = response_data.get("message", {})
    return message.get("content", "")


def read_text_file(file_path: Path, max_chars: int) -> str:
    try:
        return file_path.read_text(encoding="utf-8", errors="ignore")[:max_chars]
    except OSError:
        return ""


def read_pdf(file_path: Path, max_chars: int) -> str:
    try:
        reader = PdfReader(str(file_path))
    except Exception:
        return ""

    parts = []
    total = 0
    for page in reader.pages[:8]:
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if not text:
            continue
        remaining = max_chars - total
        if remaining <= 0:
            break
        chunk = text[:remaining]
        parts.append(chunk)
        total += len(chunk)
    return "\n".join(parts).strip()


def read_docx(file_path: Path, max_chars: int) -> str:
    try:
        document = Document(str(file_path))
    except Exception:
        return ""

    parts = []
    total = 0
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        remaining = max_chars - total
        if remaining <= 0:
            break
        chunk = text[:remaining]
        parts.append(chunk)
        total += len(chunk)
    return "\n".join(parts).strip()


def read_pptx(file_path: Path, max_chars: int) -> str:
    try:
        presentation = Presentation(str(file_path))
    except Exception:
        return ""

    parts = []
    total = 0
    for slide in presentation.slides[:10]:
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if not text:
                continue
            remaining = max_chars - total
            if remaining <= 0:
                break
            chunk = text[:remaining]
            parts.append(chunk)
            total += len(chunk)
        if total >= max_chars:
            break
    return "\n".join(parts).strip()


def read_xlsx(file_path: Path, max_chars: int) -> str:
    try:
        workbook = load_workbook(str(file_path), read_only=True, data_only=True)
    except Exception:
        return ""

    parts = []
    total = 0
    for sheet in workbook.worksheets[:4]:
        line = f"[Sheet] {sheet.title}"
        parts.append(line)
        total += len(line)
        for row in sheet.iter_rows(min_row=1, max_row=20, values_only=True):
            cells = [str(value).strip() for value in row if value not in (None, "")]
            if not cells:
                continue
            line = " | ".join(cells)
            remaining = max_chars - total
            if remaining <= 0:
                workbook.close()
                return "\n".join(parts).strip()
            chunk = line[:remaining]
            parts.append(chunk)
            total += len(chunk)
    workbook.close()
    return "\n".join(parts).strip()


def read_zip_listing(file_path: Path, max_chars: int) -> str:
    try:
        with zipfile.ZipFile(file_path) as archive:
            names = archive.namelist()[:40]
    except Exception:
        return ""

    return "\n".join(names)[:max_chars]


def read_via_textutil(file_path: Path, max_chars: int) -> str:
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(file_path)],
            capture_output=True,
            text=True,
            check=False,
        )
    except OSError:
        return ""

    if result.returncode != 0:
        return ""

    text = result.stdout.strip()
    if text.startswith("%PDF-"):
        return ""
    return text[:max_chars]


def extract_text_preview(file_path: Path, max_chars: int) -> str:
    suffix = file_path.suffix.lower()
    if suffix in {".txt", ".py", ".m", ".html", ".csv"}:
        return read_text_file(file_path, max_chars)

    if suffix == ".pdf":
        return read_pdf(file_path, max_chars)

    if suffix == ".docx":
        return read_docx(file_path, max_chars)

    if suffix in {".doc", ".rtf", ".odt", ".pages"}:
        return read_via_textutil(file_path, max_chars)

    if suffix == ".pptx":
        return read_pptx(file_path, max_chars)

    if suffix in {".xlsx", ".xls"}:
        return read_xlsx(file_path, max_chars)

    if suffix == ".zip":
        return read_zip_listing(file_path, max_chars)

    return ""


def should_use_ai(file_path: Path, preview: str, aggressive: bool = False) -> bool:
    if aggressive:
        if file_path.suffix.lower() not in VALID_EXTENSIONS:
            return False
    elif file_path.suffix.lower() not in AI_REVIEWABLE_EXTENSIONS:
        return False

    if should_force_unsorted(file_path.name):
        return False

    if aggressive:
        return True

    if has_academic_hint(file_path.name):
        return True

    if preview and has_academic_hint(preview):
        return True

    return False


def should_use_ai_for_text(
    name_or_hint: str,
    preview: str,
    is_dir: bool = False,
    aggressive: bool = False,
) -> bool:
    normalized_name = normalize(name_or_hint)

    if should_force_unsorted(normalized_name):
        return False

    if aggressive:
        return True

    if is_dir:
        if has_academic_hint(normalized_name):
            return True
        if preview and has_academic_hint(preview):
            return True
        return False

    if has_academic_hint(normalized_name):
        return True

    if preview and has_academic_hint(preview):
        return True

    return False


def build_ollama_payload(
    item_name: str,
    item_kind: str,
    item_suffix: str,
    preview: str,
    subjects: list[str],
    model: str,
    aggressive: bool = False,
) -> dict:
    prompt = (
        f"Classifica este {item_kind} de Downloads.\n"
        f"Nome: {item_name}\n"
        f"Extensao: {item_suffix}\n"
        "Pastas permitidas: "
        + ", ".join(subjects)
        + "\n"
        "Se nao houver confianca suficiente, responde com Por_Organizar.\n"
        "Nao inventes pastas novas.\n"
        "So classifiques numa cadeira se houver evidencia explicita no nome ou no texto.\n"
        "Palavras genericas como exemplo, ficheiro, printable, pdf, tickets ou declaracao nao chegam.\n"
        "Instaladores, curriculos, bilhetes, capturas de ecra e ficheiros pessoais devem ficar em Por_Organizar.\n"
        "Devolve apenas JSON com: subject, confidence, reason.\n"
    )
    if aggressive:
        prompt += (
            "Modo agressivo: tenta classificar mesmo com pistas fracas, mas nunca atribuas uma cadeira se as pistas "
            "forem genericas, pessoais ou de software sem contexto academico.\n"
        )
    if preview:
        prompt += f"Excerto do conteudo:\n{preview}\n"

    return {
        "model": model,
        "stream": False,
        "messages": [
            {
                "role": "user",
                "content": prompt,
            }
        ],
        "format": {
            "type": "object",
            "properties": {
                "subject": {
                    "type": "string",
                    "enum": subjects,
                },
                "confidence": {
                    "type": "integer",
                    "minimum": 0,
                    "maximum": 100,
                },
                "reason": {
                    "type": "string",
                },
            },
            "required": ["subject", "confidence", "reason"],
        },
        "options": {
            "temperature": 0,
        },
    }


def classify_with_ai(
    file_path: Path,
    model: str,
    min_confidence: int,
    max_text_chars: int,
    cache: dict,
    ollama_url: str,
    aggressive: bool = False,
) -> tuple[str | None, int, str]:
    subjects = list(SUBJECT_DESTINATIONS.keys()) + ["Por_Organizar"]
    cache_key = build_cache_key(file_path, model, aggressive=aggressive)
    cached = cache.get(cache_key)
    if cached:
        return cached["subject"], cached["confidence"], cached["reason"]

    preview = extract_text_preview(file_path, max_text_chars)
    if not should_use_ai(file_path, preview, aggressive=aggressive):
        return None, 0, "sem pistas suficientes para classificacao segura"

    payload = build_ollama_payload(
        file_path.name,
        "ficheiro",
        file_path.suffix.lower(),
        preview,
        subjects,
        model,
        aggressive=aggressive,
    )
    request = urllib.request.Request(
        ollama_url,
        data=json.dumps(payload).encode("utf-8"),
        headers={
            "Content-Type": "application/json",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(request, timeout=90) as response:
            response_data = json.loads(response.read().decode("utf-8"))
    except urllib.error.HTTPError as exc:
        error_body = exc.read().decode("utf-8", errors="replace")
        raise RuntimeError(f"Erro HTTP do Ollama ({exc.code}): {error_body}") from exc
    except urllib.error.URLError as exc:
        raise RuntimeError(
            "Erro ao contactar o Ollama. Verifique se a app/servidor Ollama está aberto."
        ) from exc

    raw_text = extract_response_text(response_data)
    if not raw_text:
        raise RuntimeError(f"Resposta do Ollama sem texto útil para {file_path.name}.")

    try:
        parsed = json.loads(raw_text)
    except json.JSONDecodeError as exc:
        raise RuntimeError(f"JSON inválido devolvido pelo Ollama para {file_path.name}: {raw_text}") from exc

    subject = parsed.get("subject", "Por_Organizar")
    confidence = int(parsed.get("confidence", 0))
    reason = str(parsed.get("reason", "")).strip()

    if subject not in subjects:
        subject = "Por_Organizar"

    if confidence < min_confidence:
        subject = "Por_Organizar"

    cache[cache_key] = {
        "subject": subject,
        "confidence": confidence,
        "reason": reason,
    }
    save_cache(cache)
    return subject, confidence, reason


def classify_directory_with_ai(
    dir_path: Path,
    model: str,
    min_confidence: int,
    max_text_chars: int,
    cache: dict,
    ollama_url: str,
    aggressive: bool = False,
) -> tuple[str | None, int, str]:
    subjects = list(SUBJECT_DESTINATIONS.keys()) + ["Por_Organizar"]
    cache_key = build_cache_key(dir_path, model, aggressive=aggressive)
    cached = cache.get(cache_key)
    if cached:
        return cached["subject"], cached["confidence"], cached["reason"]

    preview = collect_directory_preview(dir_path, max_text_chars)
    descriptor = f"{dir_path.name} {preview}".strip()
    if not should_use_ai_for_text(descriptor, preview, is_dir=True, aggressive=aggressive):
        return None, 0, "sem pistas suficientes para classificacao segura"

    payload = build_ollama_payload(
        dir_path.name,
        "pasta",
        "(sem extensao)",
        preview,
        subjects,
        model,
        aggressive=aggressive,
    )
    request = urllib.request.Request(
        ollama_url,
        data=json.dumps(payload).encode("utf-8"),
        headers={"Content-Type": "application/json"},
        method="POST",
    )

    try:
        with urllib.request.urlopen(request, timeout=90) as response:
            response_data = json.loads(response.read().decode("utf-8"))
    except urllib.error.HTTPError as exc:
        error_body = exc.read().decode("utf-8", errors="replace")
        raise RuntimeError(f"Erro HTTP do Ollama ({exc.code}): {error_body}") from exc
    except urllib.error.URLError as exc:
        raise RuntimeError(
            "Erro ao contactar o Ollama. Verifique se a app/servidor Ollama está aberto."
        ) from exc

    raw_text = extract_response_text(response_data)
    if not raw_text:
        raise RuntimeError(f"Resposta do Ollama sem texto útil para {dir_path.name}.")

    try:
        parsed = json.loads(raw_text)
    except json.JSONDecodeError as exc:
        raise RuntimeError(f"JSON inválido devolvido pelo Ollama para {dir_path.name}: {raw_text}") from exc

    subject = parsed.get("subject", "Por_Organizar")
    confidence = int(parsed.get("confidence", 0))
    reason = str(parsed.get("reason", "")).strip()

    if subject not in subjects or confidence < min_confidence:
        subject = "Por_Organizar"

    cache[cache_key] = {
        "subject": subject,
        "confidence": confidence,
        "reason": reason,
    }
    save_cache(cache)
    return subject, confidence, reason


def organize_downloads(
    dry_run: bool = False,
    use_ai: bool = False,
    ai_aggressive: bool = False,
    move_unsorted: bool = False,
    process_directories: bool = True,
    model: str = DEFAULT_MODEL,
    min_confidence: int = DEFAULT_MIN_CONFIDENCE,
    max_text_chars: int = DEFAULT_MAX_TEXT_CHARS,
    limit: int | None = None,
    ollama_url: str = OLLAMA_API_URL,
) -> None:
    if not dry_run:
        ensure_folder(BASE_FOLDER)
        ensure_folder(DEST_ROOT)
        ensure_folder(UNSORTED_FOLDER)

    cache = load_cache() if use_ai else {}
    moved_count = 0
    unsorted_count = 0
    kept_count = 0
    ai_count = 0
    processed = 0

    for item in sorted(BASE_FOLDER.iterdir(), key=lambda path: path.name.lower()):
        if limit is not None and processed >= limit:
            break

        if item.name.startswith("."):
            continue

        processed += 1
        subject = None
        source = "rules"
        confidence = 100

        if item.is_file():
            if item.suffix.lower() not in VALID_EXTENSIONS:
                processed -= 1
                continue

            subject = classify_with_rules(item.name)
            if subject is None and use_ai:
                subject, confidence, _reason = classify_with_ai(
                    item,
                    model=model,
                    min_confidence=min_confidence,
                    max_text_chars=max_text_chars,
                    cache=cache,
                    ollama_url=ollama_url,
                    aggressive=ai_aggressive,
                )
                source = "ai"
                ai_count += 1
                if subject == "Por_Organizar":
                    subject = None
        elif item.is_dir() and process_directories:
            if should_skip_directory(item):
                continue

            hint_text = f"{item.name} {list_directory_hints(item)}".strip()
            subject = classify_with_rules(hint_text)
            if subject is None and use_ai:
                subject, confidence, _reason = classify_directory_with_ai(
                    item,
                    model=model,
                    min_confidence=min_confidence,
                    max_text_chars=max_text_chars,
                    cache=cache,
                    ollama_url=ollama_url,
                    aggressive=ai_aggressive,
                )
                source = "ai"
                ai_count += 1
                if subject == "Por_Organizar":
                    subject = None
        else:
            processed -= 1
            continue

        destination = destination_for_subject(subject)

        if destination is not None:
            final_path = safe_move(item, destination, dry_run=dry_run)
            if source == "ai":
                print(f"[AI:{confidence:02d}] {item.name} -> {display_destination(final_path.parent)}")
            else:
                print(f"[OK] {item.name} -> {display_destination(final_path.parent)}")
            moved_count += 1
        else:
            if move_unsorted:
                final_path = safe_move(item, UNSORTED_FOLDER, dry_run=dry_run)
                if source == "ai":
                    print(f"[AI:?{confidence:02d}] {item.name} -> {display_destination(final_path.parent)}")
                else:
                    print(f"[?] {item.name} -> {display_destination(final_path.parent)}")
                unsorted_count += 1
            else:
                if source == "ai":
                    print(f"[KEEP:{confidence:02d}] {item.name} -> Downloads")
                else:
                    print(f"[KEEP] {item.name} -> Downloads")
                kept_count += 1

    print(f"\nPasta organizada: {BASE_FOLDER}")
    print(f"Destino: {DEST_ROOT}")
    print(f"Modo simulacao: {'sim' if dry_run else 'nao'}")
    print(f"Modo IA: {'sim' if use_ai else 'nao'}")
    print(f"Modo IA agressivo: {'sim' if ai_aggressive else 'nao'}")
    print(f"Processar pastas: {'sim' if process_directories else 'nao'}")
    print(f"Mover nao classificados: {'sim' if move_unsorted else 'nao'}")
    if use_ai:
        print(f"Modelo: {model}")
        print(f"Chamadas IA: {ai_count}")
        print(f"Cache: {CACHE_FILE}")
    print(f"Movidos para cadeiras: {moved_count}")
    print(f"Movidos para Por_Organizar: {unsorted_count}")
    print(f"Deixados em Downloads: {kept_count}")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Organiza ficheiros da pasta Downloads com regras locais e opcionalmente com Ollama local."
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Mostra o que seria movido sem alterar ficheiros.",
    )
    parser.add_argument(
        "--use-ai",
        action="store_true",
        help="Usa o Ollama local para classificar ficheiros nao reconhecidos pelas regras.",
    )
    parser.add_argument(
        "--ai-aggressive",
        action="store_true",
        help="Usa IA em quase todos os ficheiros e pastas validos, mesmo com pistas fracas.",
    )
    parser.add_argument(
        "--move-unsorted",
        action="store_true",
        help="Move itens nao classificados para Por_Organizar. Por omissao ficam em Downloads.",
    )
    parser.add_argument(
        "--no-directories",
        action="store_true",
        help="Nao processa pastas, apenas ficheiros.",
    )
    parser.add_argument(
        "--model",
        default=DEFAULT_MODEL,
        help=f"Modelo Ollama a usar com --use-ai. Omissao: {DEFAULT_MODEL}.",
    )
    parser.add_argument(
        "--min-confidence",
        type=int,
        default=DEFAULT_MIN_CONFIDENCE,
        help=f"Confianca minima da IA para mover ficheiros. Omissao: {DEFAULT_MIN_CONFIDENCE}.",
    )
    parser.add_argument(
        "--max-text-chars",
        type=int,
        default=DEFAULT_MAX_TEXT_CHARS,
        help=f"Numero maximo de caracteres de conteudo local enviados ao modelo. Omissao: {DEFAULT_MAX_TEXT_CHARS}.",
    )
    parser.add_argument(
        "--limit",
        type=int,
        default=None,
        help="Processa apenas os primeiros N ficheiros. Útil para testes.",
    )
    parser.add_argument(
        "--ollama-url",
        default=os.environ.get("OLLAMA_API_URL", OLLAMA_API_URL),
        help=f"URL do endpoint chat do Ollama. Omissao: {OLLAMA_API_URL}.",
    )
    return parser.parse_args()


if __name__ == "__main__":
    args = parse_args()
    organize_downloads(
        dry_run=args.dry_run,
        use_ai=args.use_ai or args.ai_aggressive,
        ai_aggressive=args.ai_aggressive,
        move_unsorted=args.move_unsorted,
        process_directories=not args.no_directories,
        model=args.model,
        min_confidence=args.min_confidence,
        max_text_chars=args.max_text_chars,
        limit=args.limit,
        ollama_url=args.ollama_url,
    )
